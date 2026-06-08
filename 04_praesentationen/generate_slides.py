"""
Mirrou Creative Studio — Abschlusspräsentation PPTX Generator (32 Slides)
Dark Luxury Design: Deep Onyx #080808 + Gold #C8A25A + Ivory #F2EFE9
exactly 8 slides per presenter (Yahya, Ralph, Olha, Denys)
Import into Google Slides: File → Import slides
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand Colors
BG_COLOR = RGBColor(0x08, 0x08, 0x08)        # Deep Onyx
SURFACE = RGBColor(0x11, 0x11, 0x13)          # Surface
GOLD = RGBColor(0xC8, 0xA2, 0x5A)            # Mirrou Gold
GOLD_LIGHT = RGBColor(0xE4, 0xC0, 0x7A)       # Gold Light
IVORY = RGBColor(0xF2, 0xEF, 0xE9)           # Ivory
BODY_COLOR = RGBColor(0xB8, 0xB4, 0xAE)      # Body Text
MUTED = RGBColor(0x6E, 0x6B, 0x66)           # Muted

# Presenter Accent Colors
PRESENTER_COLORS = {
    "Yahya": GOLD,
    "Ralph": RGBColor(0xD8, 0xD3, 0xCB),
    "Olha": RGBColor(0xE0, 0x7A, 0x5F),
    "Denys": RGBColor(0x4A, 0x90, 0xE2)
}

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=BG_COLOR):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gold_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_name='Inter', color=IVORY, bold=False, alignment=PP_ALIGN.LEFT,
                 italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    return txBox

def add_speaker_badge(slide, name, color):
    # Small speaker badge in top right
    add_text_box(slide, Inches(10.0), Inches(0.5), Inches(2.5), Inches(0.4),
                 f"SPRECHER: {name.upper()}", font_size=10, font_name='JetBrains Mono',
                 color=color, bold=True, alignment=PP_ALIGN.RIGHT)

def make_base_slide(num, title, presenter, notes_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide)
    
    # Speaker Notes
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text
        
    # Presenter Color
    p_color = PRESENTER_COLORS.get(presenter, GOLD)
    
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.4),
                 f"{num} · {title.upper()}", font_size=11,
                 font_name='JetBrains Mono', color=p_color, bold=True)
                 
    add_speaker_badge(slide, presenter, p_color)
    return slide

# 32 Slide Data Definitions
slide_data = [
    # YAHYA (1-6)
    {
        "num": "01",
        "title": "Titel & Begruessung",
        "presenter": "Yahya",
        "notes": "Guten Tag und herzlich willkommen zur Abschlusspräsentation von Mirrou Creative Studio. Mein Name ist Yahya Yildirim. Gemeinsam mit unserer Gründerin Olha Yevtushenko und meinen Partnern Denys Demyanyshyn und Ralph Kindermann zeigen wir Ihnen heute den Aufbau einer KI-integrierten Performance Creative Agency. Wir demonstrieren, wie wir über unsere Frontier-Firm-Architektur, einen hybriden Produktions-Workflow und proaktive Compliance eine neue Kategorie besetzen: Creative Performance für Beauty, Health und Lifestyle. Lassen Sie uns einen Blick darauf werfen, wie wir uns als Team aufgeteilt und dieses Projekt realisiert haben.",
        "content_type": "title_slide",
        "header": "Aufbau einer KI-integrierten\nPerformance Creative Agency",
        "sub": "Frontier Firm Architektur · Hybrid Production · EU AI Act-Compliance"
    },
    {
        "num": "02",
        "title": "Aufwand & Mission",
        "presenter": "Yahya",
        "notes": "Um ein verteiltes Team über zwei Städte hinweg fehlerfrei zu orchestrieren, arbeiten wir nach dem Prinzip einer Frontier Firm. In Hamburg steuert Olha als unsere Gründerin und Creative Director die gesamte visuelle Marken-Identität und Produktion (~120h), unterstützt von Denys für Campaign Management und Marketing-KI (~80h). In Berlin leitet Ralph das CRM und Lifecycle-Marketing (~60h) und ich verantworte das Growth- und Inbound-Marketing (~150h). Insgesamt haben wir ca. 410 Stunden realen Build investiert, um Mirrou und unsere Infrastruktur live an den Start zu bringen. Doch warum brennt der Markt gerade ab?",
        "content_type": "metrics",
        "header": "4 Koepfe. 410 Stunden. Reale Wirkung.",
        "metrics": [
            ("Olha Yevtushenko", "Gründerin & Creative Director", "120 Std."),
            ("Denys Demyanyshyn", "Campaign Manager & Marketing-KI", "80 Std."),
            ("Ralph Kindermann", "CRM & Lifecycle Marketing", "60 Std."),
            ("Yahya Yildirim", "Growth & Inbound (Projektlead)", "150 Std.")
        ]
    },
    {
        "num": "03",
        "title": "Marktproblem: Fatigue",
        "presenter": "Yahya",
        "notes": "Der Grund für Mirrou liegt in einem lautlosen Killer im Paid-Social-Bereich: Creative Fatigue. Beauty-, Health- und Lifestyle-Marken im DACH-Raum stoßen ab einem Monats-Spend von 20.000 Euro an ein unsichtbares Limit. Daten aus dem Meta-Ecosystem zeigen, dass Werbemittel durch die hohe Reizüberflutung bereits nach 14 bis 28 Tagen ausbrennen. Das führt zu einem dramatischen Klickraten-Einbruch von 35 bis 55 Prozent und einem Anstieg der CPC-Kosten um bis zu 70 Prozent. Klassische Kreativagenturen reagieren darauf zu langsam, reine Performance-Agenturen oft zu seelenlos.",
        "content_type": "two_col",
        "header": "Das Marktproblem: Creative Fatigue",
        "col1_title": "DER LAUTLOSE KILLER IM PAID SOCIAL",
        "col1_text": "Werbemittel brennen bei D2C-Brands nach 14-28 Tagen komplett aus. Budgets verbrennen durch sinkende Klickraten und steigende Kosten.",
        "col2_title": "BENCHMARKS (FATIGUE-ONSET)",
        "col2_text": "• CTR Meta Feed: -35% bis -55% Einbruch\n• CPC Meta: +45% bis +70% Anstieg\n• Halbwertszeit: 14-28 Tage (ab 20k Euro/Monat Spend)"
    },
    {
        "num": "04",
        "title": "Wettbewerbs-Arenen",
        "presenter": "Yahya",
        "notes": "Unsere Positionierung verortet sich im Schnittpunkt dreier Markt-Arenen. Klassische Agenturen vernachlässigen oft die Conversion-Daten; reine Performance-Klitschen opfern die visuelle Ästhetik. Mirrou löst dieses Dilemma durch die algorithmische Seele: Wir kombinieren Premium-Design mit einer hochgradig automatisierten Testing-Engine, um Creative Fatigue systematisch zu eliminieren und den Return on Ad Spend nachhaltig abzusichern.",
        "content_type": "two_col",
        "header": "Drei Arenen. Eine leere Ecke.",
        "col1_title": "DIE WETTBEWERBS-LANDSCHAFT",
        "col1_text": "• Boutique-Studios: Hohe Ästhetik, kein Performance-System, extrem langsam (4-8 Wochen).\n• Performance-Agenturen: Starke Daten, aber visuell austauschbar.\n• In-House-Teams: Hohe Fixkosten, langsame Skalierung.",
        "col2_title": "DIE BESETZTE NICHE: MIRROU",
        "col2_text": "Mirrou vereint die emotionale Seele von Premium-Design mit der kalten Datenlogik von A/B-Testing und künstlicher Intelligenz in der Beauty/Health Nische."
    },
    {
        "num": "05",
        "title": "Service-Pakete & Retainer",
        "presenter": "Yahya",
        "notes": "Unser Angebot richtet sich an D2C-Brands mit 10.000 bis 150.000 Euro monatlichem Ad-Spend, die keine Kompromisse machen wollen. Wir bieten drei standardisierte, direkt buchbare Service-Pakete auf der Website: Einmalige Shootings im Paket ‚E-Commerce & Catalog‘ zur Katalog-Pflege, das ‚Social Media & Advertising‘ Paket für gezielte Kampagnen-Launches und unseren ‚Creative Retainer‘ in den Tiers S, M und L für eine fortlaufende Creative-Engine im Monatsrhythmus.",
        "content_type": "three_col",
        "header": "Drei Wege, mit uns zu arbeiten",
        "col1": ("PAKET 01", "E-Commerce & Catalog", "Online-Schaufenster / Reinzeichnung\nPreis: 1.500 - 3.000 Euro"),
        "col2": ("PAKET 02", "Social & Advertising", "Hypothesen-getestete Video/Still Ad Sets\nPreis: 2.000 - 5.000 Euro"),
        "col3": ("PAKET 03", "Creative Retainer", "Monatliche Creative Engine (Tiers S/M/L)\nPreis: 2.000 - 15.000 Euro/Monat")
    },
    {
        "num": "06",
        "title": "Frontier Firm 4-Layer",
        "presenter": "Yahya",
        "notes": "Um diesen kontinuierlichen Output ohne massiven Personalüberhang zu liefern, haben wir eine 4-Schichten-Infrastruktur implementiert: Den Intelligence Layer über Perplexity Spaces, den Production Layer über Adobe und Midjourney, den Infrastructure Layer auf GCP und den Performance Layer. Jede Schicht wird über ein automatisiertes GitHub-Governance-OS auditiert und versioniert.",
        "content_type": "two_col",
        "header": "Frontier Firm: 4 Menschen wie 20.",
        "col1_title": "DIGITALE SYSTEM-SCHICHTEN",
        "col1_text": "1. Intelligence Layer: Perplexity Spaces & Claude (Strategie)\n2. Production Layer: Midjourney & Adobe Creative Cloud\n3. Infrastructure Layer: Google Cloud Run & MCP Connectors\n4. Performance Layer: Meta Ads Manager, GA4 & Looker",
        "col2_title": "GITHUB GOVERNANCE",
        "col2_text": "Das gesamte System wird in einem versionierten GitHub-Repository als Single Source of Truth verwaltet und versioniert. Voller Audit-Trail über Commit-History."
    },

    # RALPH (7-8)
    {
        "num": "07",
        "title": "Perfect-Twin-Prinzip",
        "presenter": "Ralph",
        "notes": "Vielen Dank, Yahya. Das Geheimnis hinter unserer operativen Effizienz liegt im Perfect-Twin-Prinzip: Jeder menschliche Spezialist wird durch ein maßgeschneidertes KI-Pendant erweitert. So übersetzt Olha ihre Vision blitzschnell in KI-Hintergründe, Denys nutzt einen Claude Data Agent zur CSV-Analyse, ich automatisiere CRM-Workflows und Yahya steuert das React-Frontend über Claude Code. Olha, wie übersetzt du diese Architektur in Ästhetik?",
        "content_type": "two_col",
        "header": "Perfect Twin: Die Mensch-KI-Kopplung",
        "col1_title": "HUMAN EXPERTISE + AI TWIN",
        "col1_text": "Jeder Spezialist nutzt ein maßgeschneidertes KI-Gegenstück zur Skalierung.\n• Olha (CD) + Midjourney/Firefly\n• Denys (Data) + Claude Data Agent",
        "col2_title": "PROZESS-INTEGRATION",
        "col2_text": "• Ralph (CRM) + HubSpot CRM Automation\n• Yahya (Tech) + Claude Code Developer Tools\nDas Ergebnis: 5x Zeiteinsparung und absolute Konsistenz im Output."
    },
    {
        "num": "08",
        "title": "Inbound Social Setup",
        "presenter": "Ralph",
        "notes": "Um das System mit Neukunden zu füttern, nutzen wir ein automatisiertes Inbound-Setup. Traffic wird über Thought Leadership auf LinkedIn und organische Visuals auf Instagram und Facebook generiert. Sobald ein Prospect die Website besucht, durchläuft er ein automatisiertes Formular-Scoring und wird direkt in unsere Pipeline geroutet. Olha, wie erwacht diese Marke visuell zum Leben?",
        "content_type": "two_col",
        "header": "Inbound-Routing: Vom Klick zum Lead",
        "col1_title": "TRAFFIC GENERATION",
        "col1_text": "• LinkedIn: B2B Thought-Leadership Artikel\n• Meta/Instagram: Organische und bezahlte Product Visuals\n• Facebook: Community & Social Proof",
        "col2_title": "ROUTING & SCORING",
        "col2_text": "Prospects landen auf mirrou.studio, füllen das Ad-Spend-Formular aus und werden basierend auf einem automatischen Qualifizierungs-Score direkt geroutet."
    },

    # OLHA (9-16)
    {
        "num": "09",
        "title": "Visuelle Philosophie",
        "presenter": "Olha",
        "notes": "Vielen Dank, Ralph. Als Creative Director & Performance Marketer ist meine Mission, den algorithmischen Rahmen mit einer Seele zu füllen. Unsere Antwort darauf ist die Synergie aus Ästhetik und Conversion über unser ‚Dark Luxury‘-System. Wir designen keinen Pixel ohne Performance-Zweck: Dunkle Beleuchtung stellt das Produkt in den Mittelpunkt, während der Text minimal und funktional bleibt, um Ladezeiten und CTR zu optimieren.",
        "content_type": "two_col",
        "header": "Visuelle Philosophie: Aesthetik + Conversion",
        "col1_title": "ALGORITHM OF SOUL",
        "col1_text": "Wir brechen mit dem typischen, hellen Beauty-Markt. Dark Luxury erzeugt visuelle Unterbrechung im Scroll-Feed.",
        "col2_title": "DESIGN RULES FOR PERFORMANCE",
        "col2_text": "• Low-Key Beleuchtung: Produkt plastisch im Fokus\n• Minimaler Textanteil: Schnelle Erfassung und Ladezeiten\n• Fokus: Makro-Texturen von Haut und Flüssigkeiten"
    },
    {
        "num": "10",
        "title": "Brand Identity Tokens",
        "presenter": "Olha",
        "notes": "Die visuelle Identität wird über programmatische Tokens gesteuert. Die Basis bildet ein sattes Deep Onyx (#080808) für visuelle Ruhe. Muted Gold (#C8A25A) setzen wir kontrolliert als Akzent ein. Kontrastiert wird das System durch ein warmes Ivory (#F2EFE9). Unsere Typografie kombiniert die luxuriöse Display-Serif Cormorant Garamond mit die präzisen serifenlosen Inter und technischen Details in JetBrains Mono.",
        "content_type": "three_col",
        "header": "Brand Identity Tokens",
        "col1": ("FARBTAPETE", "Deep Onyx", "#080808 / #0A0A0A\nErzeugt edle visuelle Ruhe"),
        "col2": ("AKZENT", "Muted Gold", "#C8A25A\nFiligran und edel"),
        "col3": ("KONTRAST", "Warm Ivory", "#F2EFE9\nOptimale Lesbarkeit und Eleganz")
    },
    {
        "num": "11",
        "title": "Brandbook & Rules",
        "presenter": "Olha",
        "notes": "Qualität und Konsistenz sichern wir durch ein 49-seitiges Brandbook. Es regelt Grid-Systeme, Bildsprache, Beleuchtungs-Vorgaben und Do's & Don'ts. Jedes Teammitglied und jeder externe Creator kann auf Basis dieser klaren Guidelines visuell einheitliche, konversionsstarke Werbemittel erstellen, ohne dass die Marke verwässert.",
        "content_type": "two_col",
        "header": "49-seitiges Brandbook fuer Konsistenz",
        "col1_title": "REGELN FÜR DIE KREATION",
        "col1_text": "• Grid-Systeme und Bildkomposition\n• Beleuchtungswinkel und Schattenkanten\n• Richtlinien für KI-Prompts und Seed-Consistency",
        "col2_title": "DO'S & DON'TS",
        "col2_text": "• Do: Produkt steht immer im echten Fotomittelpunkt\n• Don't: Überladene Banner, bunte Overlays, künstliche Gesichter"
    },
    {
        "num": "12",
        "title": "Demo Case 1: Skincare",
        "presenter": "Olha",
        "notes": "Um unsere Bandbreite zu beweisen, haben wir vier Demo-Cases entwickelt. Unser erster Case, ‚Luminous Aura‘, zeigt ein Premium-Serum auf weißem Marmor mit weichem Goldstaub – optimiert für klassische Conversions im Meta-Kanal. Das Visual kontrastiert das Produkt und führte in A/B-Tests zu einer CTR-Steigerung um 82 Prozent.",
        "content_type": "two_col",
        "header": "Demo Case 1: Luminous Aura",
        "col1_title": "DAS PRODUKT IM FOKUS",
        "col1_text": "Premium-Serum auf weißem Marmor. Minimales Design, kontrastierende Details. Weicher Goldstaub symbolisiert Luxus.",
        "col2_title": "PERFORMANCE METRIK (SIMULIERT)",
        "col2_text": "• Kanal: Meta Feed (Instagram/Facebook)\n• A/B-Test: +82% CTR Steigerung vs. Standard-Studio\n• Conversion Rate: Signifikant erhöht durch Premium-Trust"
    },
    {
        "num": "13",
        "title": "Demo Cases 2 & 3: Motion",
        "presenter": "Olha",
        "notes": "Unsere Cases 2 und 3 zeigen weitere visuelle Nischen: ‚Vitality Pulse‘ setzt auf kinetische Energie, Wassertropfen und harte Schatten für TikTok, optimiert für schnelle Hook-Variationen über die ersten 3 Sekunden. ‚Essence Drift‘ nutzt Nebel und Glasspiegelungen für emotionales Storytelling im Pinterest- und Instagram-Reels-Kanal.",
        "content_type": "two_col",
        "header": "Demo Cases: Motion & Emotion",
        "col1_title": "CASE 2: VITALITY PULSE (TIKTOK)",
        "col1_text": "Kinetische Energie, Wassertropfen, harte Schatten. Optimiert auf die ersten 3 Sekunden (Thumb-Stop-Hooks).",
        "col2_title": "CASE 3: ESSENCE DRIFT (REELS/PINTEREST)",
        "col2_text": "Nebel, Glasspiegelungen, emotionale Tonalität. Fokus auf Fragrance-Storytelling und Customer Desire."
    },
    {
        "num": "14",
        "title": "Demo Case 4: KI Biotech",
        "presenter": "Olha",
        "notes": "Unser vierter Case, ‚Neural Glow‘, visualisiert biotech-orientierte Skincare über fluoreszierende Schaltkreise und ist vollständig KI-generiert. Dieser Case beweist, dass wir synthetische Bildwelten fotorealistisch beherrschen. Gleichzeitig sind alle Bilddaten mit C2PA-Metadaten versehen, was uns absolute Rechtssicherheit nach dem kommenden EU AI Act garantiert.",
        "content_type": "two_col",
        "header": "Demo Case 4: Neural Glow (100% KI)",
        "col1_title": "SYNTHETISCHE WELTEN",
        "col1_text": "Biotech-orientierte Skincare visualisiert über fluoreszierende Schaltkreise. Vollständig digital gerendert.",
        "col2_title": "EU AI ACT READY",
        "col2_text": "• Kryptografische C2PA-Metadaten in der Datei eingebettet\n• Rechtssichere Transparenzkennzeichnung nach Art. 50"
    },
    {
        "num": "15",
        "title": "Hybrid Workflow Flow",
        "presenter": "Olha",
        "notes": "Unser Kernprozess ist die Hybrid Production. Wir trennen das physische Produkt von seiner Umgebung: Das Produkt wird in unserem Hamburger Studio real fotografiert. Die komplexen Hintergründe und Lichtstimmungen werden jedoch digital per Midjourney gerendert. Das spart uns den aufwendigen analogen Set-Bau und senkt die Kosten drastisch.",
        "content_type": "two_col",
        "header": "Hybrid Production Workflow",
        "col1_title": "01. ECHTE FOTOGRAFIE (HAMBURG)",
        "col1_text": "Das physische Produkt, Texturen und Flakon werden real im Studio fotografiert. Keine KI-Verfälschung des Kernprodukts.",
        "col2_title": "02. DIGITAL COMPOSING (KI)",
        "col2_text": "Hintergründe, Lichtstimmungen und Reflektionen werden per Midjourney/Firefly generiert und fotorealistisch verschmolzen."
    },
    {
        "num": "16",
        "title": "Workflow Effizienz-Beweis",
        "presenter": "Olha",
        "notes": "Die wirtschaftliche Hebelwirkung der Hybrid Production is enorm. Während ein klassisches Studio-Shooting inklusive Set-Bau, Vor-Ort-Produktion und Post-Production circa 144 Stunden (6 Tage) in Anspruch nimmt, schrumpft unser Prozess auf 4 Stunden digitales Composing. Das ist eine Zeitersparnis von über 97 Prozent. Denys, wie misst du den Erfolg dieser Assets?",
        "content_type": "two_col",
        "header": "Hybrid Production Effizienz-Beweis",
        "col1_title": "KLASSISCHES STUDIO-SHOOTING",
        "col1_text": "• Set-Bau, Studio-Miete, Fotografen, Retusche\n• Aufwand: ca. 144 Std. (6 Tage)\n• Starre Setups, keine schnellen Änderungen möglich",
        "col2_title": "MIRROU HYBRID SHOOTING",
        "col2_text": "• Einmaliges Produkt-Shooting + KI-Composing\n• Aufwand: ca. 4 Std. (Digitale Pipeline)\n• Zeit- und Kostenersparnis: über 97%\n• Skalierbarkeit: Unbegrenzt viele Hintergründe"
    },

    # DENYS (17-24)
    {
        "num": "17",
        "title": "Performance & QFC",
        "presenter": "Denys",
        "notes": "Danke, Olha. Wo deine visuelle Exzellenz aufhört, beginnt mein Messsystem. Wir optimieren Kampagnen nicht nur nach historischen Klicks. Über Google Marketing Live 2026 integrieren wir Qualified Future Conversions (QFC). Diese Gemini-gestützte Metrik verknüpft frühe Nutzersignale wie Video-Views und Branded Search mit prognostizierten Käufen bis zu 6 Monate im Voraus.",
        "content_type": "two_col",
        "header": "Performance: Qualified Future Conversions",
        "col1_title": "BEYOND HISTORICAL DATA",
        "col1_text": "Klassische Performance-Optimierung blickt zurück. Mirrou blickt nach vorn über GML 2026 QFC-Modellierung.",
        "col2_title": "QUALIFIED FUTURE CONVERSIONS (QFC)",
        "col2_text": "Gemini-gestützte Predictive Engine verknüpft frühe Nutzersignale (z. B. Video-Views) mit zukünftigen Käufen (bis zu 6 Monate)."
    },
    {
        "num": "18",
        "title": "Mediabudget Donut Split",
        "presenter": "Denys",
        "notes": "Unser Mediabudget verteilen wir strategisch auf 5 Kanäle, wie Sie auf Slide 18 an dem Budget-Donut sehen können: 40 Prozent fließen in Meta Ads für Awareness über Reels-Videos und Stills. 25 Prozent gehen in TikTok Ads zur Testung schneller Hooks. 20 Prozent investieren wir in Google Search zur Conversion-Sicherung. 10 Prozent steuern das E-Mail-Marketing und 5 Prozent nutzen wir auf LinkedIn für B2B-Trust.",
        "content_type": "two_col",
        "header": "Mediabudget: 5 Kanaele im Fokus",
        "col1_title": "STRATEGISCHE BUDGETVERTEILUNG",
        "col1_text": "• Meta Ads: 40% (Kombination Reels & Still Ads)\n• TikTok Ads: 25% (Kinetische Video-Hooks)\n• Google Search: 20% (Conversion-Sicherung)",
        "col2_title": "RETENTION & TRUST",
        "col2_text": "• Email Marketing: 10% (Reaktivierung, CLV)\n• LinkedIn: 5% (B2B Authority & Trust)"
    },
    {
        "num": "19",
        "title": "Leitmetriken & Kanäle",
        "presenter": "Denys",
        "notes": "Jeder Kanal erfordert eine eigene Leitmetrik: Meta optimiert nach der Klickrate (CTR), TikTok nach der Thumb-Stop Rate, um sicherzustellen, dass User nicht wegschreiben. Google Search wird nach ROAS und CPA gesteuert, während der E-Mail-Kanal zur Reaktivierung und zur Steigerung des Customer Lifetime Value (CLV) beiträgt.",
        "content_type": "two_col",
        "header": "Kanalspezifische Leitmetriken",
        "col1_title": "META & TIKTOK LEITMETRIKEN",
        "col1_text": "• Meta: Klickrate (CTR) — Ziel >1.5%\n• TikTok: Thumb-Stop Rate (3s View) — Ziel >35%",
        "col2_title": "SEARCH & RETENTION",
        "col2_text": "• Google Search: Return on Ad Spend (ROAS) & CPA\n• Email Marketing: Conversion Rate & Customer Lifetime Value (CLV)"
    },
    {
        "num": "20",
        "title": "Daten-Feedback Loop",
        "presenter": "Denys",
        "notes": "Zur Optimierung nutzen wir den 5-Schritt-Algorithmus: Er startet mit dem Creative Audit zur Fatigue-Erkennung. Daraufhin formulieren wir 3 Performance-Hypothesen, die Olha im Hybrid-Workflow umsetzt. Wir testen diese gezielt mit minimalem Budget über 5 Tage in einer isolierten Test-Kampagne. Erst nach statistischer Signifikanz skalieren wir die Gewinner.",
        "content_type": "two_col",
        "header": "Datenkreislauf: 5-Schritt-Algorithmus",
        "col1_title": "VOM AUDIT ZUR SKALIERUNG",
        "col1_text": "1. Creative Audit: Erkennung von Fatigue und CPC-Spikes\n2. Hypothesen: Formulierung von 3 Test-Winkeln\n3. Umsetzung: Hybrid Production durch Olha",
        "col2_title": "TESTEN & LERNEN",
        "col2_text": "4. Performance Layer: Isolierter A/B-Test (5-Tage)\n5. Feedback Loop: CTR-Daten fließen zurück in Creative Engine"
    },
    {
        "num": "21",
        "title": "Testing Matrix Winkel",
        "presenter": "Denys",
        "notes": "Die systematische Generierung von Hypothesen erfolgt über unsere Creative Testing Matrix. Wir testen Creatives in vier Dimensionen: Craft (Inhaltsstoffe), Data (Studien und Wirksamkeits-Belege), Luxury (Status und Ästhetik) und Results (sichtbare Vorher-Nachher Effekte), um den stärksten Kaufanreiz der Zielgruppe zu isolieren.",
        "content_type": "two_col",
        "header": "Creative Testing Matrix",
        "col1_title": "TEST-WINKEL 1 & 2",
        "col1_text": "• Craft (Inhaltsstoffe): Fokus auf Makro-Produktaufnahmen\n• Data (Wissenschaft): Einblendung von Studien & Nachweisen",
        "col2_title": "TEST-WINKEL 3 & 4",
        "col2_text": "• Luxury (Status/Brand): Dunkles, cleanes Premium-Styling\n• Results (Ergebnis): Sichtbarer Vorher-Nachher Effekt"
    },
    {
        "num": "22",
        "title": "25-Tage Testplan Sprints",
        "presenter": "Denys",
        "notes": "Ein Kampagnen-Launch folgt einem strikten 25-Tage-Testplan: In Sprint 1 (Tag 1–8) bauen wir das Pixel-Tracking auf und erfassen die Baseline. In Sprint 2 (Tag 9–16) starten wir die A/B-Tests für Hooks und Landingpages. In Sprint 3 (Tag 17–25) skalieren wir die Gewinner. Feste Regeln steuern das Media Buying und eliminieren Emotionen.",
        "content_type": "three_col",
        "header": "25-Tage Testplan Sprints",
        "col1": ("SPRINT 1 (TAG 1-8)", "Baseline", "Pixel-Einrichtung\nTargeting Definition\nErfassung Benchmark CTR"),
        "col2": ("SPRINT 2 (TAG 9-16)", "A/B Testing", "Start A/B-Tests\n14-Tage Hook-Iterationen\nLandingpage Optimierung"),
        "col3": ("SPRINT 3 (TAG 17-25)", "Scaling", "Budget-Allokation auf Gewinner-Ads\nErschließung Lookalike-Audiences")
    },
    {
        "num": "23",
        "title": "Automated Rules Engine",
        "presenter": "Denys",
        "notes": "Um menschliche Verzögerungen beim Media-Buying auszuschließen, nutzen wir automatisierte Rules-Engines: Liegt die CTR einer Ad nach 3 Tagen unter 1,5 Prozent, wird sie automatisch pausiert und der Hook getauscht. Liegt der CPL unter 12 Euro, erhöht das System das Budget automatisch um 25 Prozent pro Tag. Ralph, wie fängt unser CRM diese Leads ab?",
        "content_type": "two_col",
        "header": "Automated Rules Engine",
        "col1_title": "IF-THIS-THEN-THAT BUDGETSTEUERUNG",
        "col1_text": "Menschliche Fehler und Verzögerungen beim Media Buying werden durch programmierte Kampagnen-Regeln eliminiert.",
        "col2_title": "REALE REGEL-ANWEISUNGEN",
        "col2_text": "• Regel 1: CTR < 1,5% nach 3 Tagen → Ad stopp & Hook-Austausch\n• Regel 2: CPL < 12 Euro → Erhöhe Budget um 25% pro Tag\n• Regel 3: Frequenz > 4.5 → Führe neue Design-Variante ein"
    },
    {
        "num": "24",
        "title": "Learning Log & Columna",
        "presenter": "Denys",
        "notes": "Jedes Testergebnis wird im Creative Learning Log archiviert, um Wissen dauerhaft zu sichern. Über OMMs Columna-Engine überwachen wir zudem kontinuierlich die Werbemittel-Adoption unserer Wettbewerber, um Markttrends und neue QFC-Modellierungen frühzeitig zu erkennen. Ralph, wie sieht unsere CRM-Infrastruktur aus?",
        "content_type": "two_col",
        "header": "Creative Learning Log & Columna",
        "col1_title": "KNOWLEDGE MANAGEMENT",
        "col1_text": "• Creative Learning Log: Archivierung aller A/B-Testergebnisse\n• Vermeidung von wiederholten Fehlern bei Ad-Iterations",
        "col2_title": "COMPETITIVE INTELLIGENCE",
        "col2_text": "OMM Columna scannt kontinuierlich Wettbewerber-Ads im Beauty-Sektor. Früherkennung von Trends und Hook-Varianten."
    },

    # RALPH (25-30)
    {
        "num": "25",
        "title": "Operations Philosophy",
        "presenter": "Ralph",
        "notes": "Vielen Dank, Denys. Qualität entsteht durch Wiederholbarkeit. Meine Aufgabe ist es, aus den ankommenden Leads ein stabiles, skalierbares Kundensystem zu formen. Das beginnt mit unserer Customer Operations Philosophy: Die vollständige Eliminierung von Kommunikations- und Delivery-Fehlern durch standardisierte digitale Abläufe.",
        "content_type": "two_col",
        "header": "Customer Operations Philosophy",
        "col1_title": "ZERO-ERROR OPERATIONS",
        "col1_text": "Standardisierung schlägt kreatives Chaos. Strukturierter HubSpot & Notion Setup minimiert Kommunikations-Reibung.",
        "col2_title": "TAKTE DER KOMMUNIKATION",
        "col2_text": "• Lead Response Zeit: Unter 15 Minuten\n• Reporting Rhythmus: Wöchentliche Looker-Updates\n• Übergabe-Prozedere der Assets: Standardisierter Drive Link"
    },
    {
        "num": "26",
        "title": "Inbound & Outbound Paths",
        "presenter": "Ralph",
        "notes": "Wir steuern Sales-Prozesse über zwei Pfade: Pfad 1 ist die Inbound-Journey, bei der ein Kunde eine Meta-Ad sieht, auf mirrou.studio landet, unser Mission Deck lädt und ein Erstgespräch bucht. Pfad 2 ist die Outbound-Journey: Ich recherchiere aktive Ads von Skincare-Brands, erstelle ein kurzes personalisiertes Video-Audit und kontaktiere den CMO auf LinkedIn.",
        "content_type": "two_col",
        "header": "Customer Journeys: Inbound & Outbound",
        "col1_title": "INBOUND JOURNEY",
        "col1_text": "Ad-Klick → mirrou.studio Landingpage → Qualifizierungsformular → Calendly Buchung → Live-Pitch & Audit",
        "col2_title": "OUTBOUND JOURNEY",
        "col2_text": "Wettbewerber Ad Library Research → Ralph erstellt 90s Video-Audit → Outreach via LinkedIn → Test-Shootings-Retainer"
    },
    {
        "num": "27",
        "title": "HubSpot CRM Pipeline",
        "presenter": "Ralph",
        "notes": "Als technisches Fundament nutzen wir HubSpot Starter als CRM. Unsere Pipeline hat 7 Stufen: Vom ersten Inbound Lead über den Intro Call und das Creative Audit bis zum Angebot und dem Status ‚Closed Won‘. Jede Deal-Stufe ist im CRM exakt dokumentiert, um jederzeit volle Übersicht über die Sales-Pipeline zu behalten.",
        "content_type": "two_col",
        "header": "HubSpot CRM Pipeline",
        "col1_title": "7 OPERATIVE PIPELINE-STUFEN",
        "col1_text": "1. 01 Inbound Lead\n2. 02 Intro Call Booked\n3. 03 Strategy/Audit\n4. 04 Proposal Sent",
        "col2_title": "DEAL STATUS ABWICKLUNG",
        "col2_text": "5. 05 Negotiation\n6. 06 Closed Won\n7. 07 Closed Lost"
    },
    {
        "num": "28",
        "title": "ICP Scoring Matrix",
        "presenter": "Ralph",
        "notes": "Um unseren Vertrieb zu fokussieren, nutzen wir einen automatisierten Qualifizierungs-Score. Sobald ein Lead das Inbound-Formular ausfüllt, vergibt das System Punkte für Ad-Spend, Branche und Erreichbarkeit. Leads mit Score 8 oder höher triggern sofort einen Slack-Alarm, sodass wir innerhalb von 15 Minuten antworten können.",
        "content_type": "two_col",
        "header": "ICP Scoring Matrix",
        "col1_title": "AUTOMATISIERTE QUALIFIZIERUNG",
        "col1_text": "• Score Punkte für monatlichen Ad-Spend\n• Score Punkte für Beauty/Supplements Nische\n• Score Punkte für B2B-Erreichbarkeit",
        "col2_title": "15-MINUTEN SLACK ALERT",
        "col2_text": "Lead-Score >= 8 → Sofortige Slack-Push-Notification an Yahya. Intro-Call-Terminierung innerhalb von 15 Min."
    },
    {
        "num": "29",
        "title": "Notion Workspace OS",
        "presenter": "Ralph",
        "notes": "Das Wissenszentrum unserer Agentur bildet Notion OS mit 4 Kernordnern: Growth, Production, Operations und Client Records. Zur Entlastung nutzen wir OMMs LYGOX-Schnittstelle: Sobald eine Deal-Stufe im CRM auf ‚Onboarding‘ wechselt, triggert LYGOX automatisch das Erstellen des Notion-Dossiers und legt die Ordnerstruktur an.",
        "content_type": "two_col",
        "header": "Notion Workspace OS",
        "col1_title": "4 ZENTRALE WISSENSORDNER",
        "col1_text": "• Growth: Outbound outreach, Sales Templates\n• Production: Brandbook, Lighting-Setups\n• Operations: Verträge, Buchhaltung, Rechnungen",
        "col2_title": "LYGOX CRM AUTOMATION",
        "col2_text": "Deal-Status 'Onboarding' → LYGOX API-Call → Notion Dossier Erstellung, Ordnerstrukturen-Anlage in Drive."
    },
    {
        "num": "30",
        "title": "SOPs & Checklisten",
        "presenter": "Ralph",
        "notes": "Jedes Teammitglied arbeitet nach festen Standard Operating Procedures (SOPs). Auf Folie 30 sehen Sie die Checklisten für SOP 01 (Lead-Handling), SOP 02 (Kunden-Onboarding) und SOP 03 (Case-Study-Erstellung). Unsere strengen Naming-Conventions garantieren, dass Daten per Perplexity-MCP-Schnittstelle in Sekundenschnelle auffindbar sind.",
        "content_type": "two_col",
        "header": "Standard Operating Procedures (SOPs)",
        "col1_title": "QUALITÄT DURCH CHECKLISTEN",
        "col1_text": "• SOP 01: Lead Handling (Terminierung & Audit)\n• SOP 02: Onboarding (Vertrag & Asset-Abfrage)\n• SOP 03: Case Study Erstellung (Zahlen-Archiv)",
        "col2_title": "NAMING CONVENTIONS",
        "col2_text": "[DATE]_[BRAND]_[PROJECT]_[VERSION]_[INITIALS]\nErlaubt sofortige Datenauffindbarkeit über MCP-Schnittstelle."
    },

    # YAHYA (31-32)
    {
        "num": "31",
        "title": "Tech Stack & GCP Deploy",
        "presenter": "Yahya",
        "notes": "Danke, Ralph. Zum Abschluss möchte ich betonen: Mirrou ist voll einsatzbereit. Unsere Plattform basiert auf React 19 mit Vite 6, TypeScript und Tailwind CSS v4, deployed auf Firebase Hosting, mit einem Python FastAPI Backend auf Google Cloud Run in Frankfurt. Die App ist unter der Staging-URL live. Und das DNS-Setup für unsere Wunschdomain mirrou.studio liegt als fertige 4-Einträge-Konfiguration für Olha bereit bei IONOS, so dass wir in wenigen Minuten auf der Hauptdomain live gehen können.",
        "content_type": "two_col",
        "header": "Tech Stack & Google Cloud Deploy",
        "col1_title": "LIGHTWEIGHT REACTION ENGINE",
        "col1_text": "• React 19 + Vite 6 + TypeScript\n• Tailwind CSS v4 + Motion animations\n• Statische Vorgenerierung (vite-react-ssg)",
        "col2_title": "DEPLOYMENT ARCHITECTURE",
        "col2_text": "• Frontend auf Firebase Hosting (CDN)\n• Backend als Python (FastAPI) Container auf Cloud Run\n• Lighthouse-Score Desktop: 100/100/100/100 · 6/6 Security Header"
    },
    {
        "num": "32",
        "title": "Manifest & SaaS Outlook",
        "presenter": "Yahya",
        "notes": "Mirrou beweist die Kraft des OMM-Substrats: ein 4-Personen-Team steuert ein hochgradig automatisiertes Studio. Nach 410 Stunden realem Build steht das System. Unser Ausblick für Phase 2: Nach dem internen Dogfooding mit Mirrou als Tenant #1 werden wir OMMs Suite – Opus Magnum, LYGOX und Columna – als B2B-SaaS für D2C-Brands öffnen. Vielen Dank für Ihre Aufmerksamkeit.",
        "content_type": "two_col",
        "header": "SaaS Outlook: Seite B",
        "col1_title": "DOGFOODING COMPLETED",
        "col1_text": "Mirrou dient als Tenant #1 zur Validierung des Opus-Magnum-Substrats (LYGOX, Columna).",
        "col2_title": "B2B RETROFIT OUTLOOK",
        "col2_text": "Öffnung der OMM Enterprise-Suite als lizenziertes B2B-SaaS für D2C-Brands ab Phase 2. mirrou.studio live."
    }
]

# Run through the 32 slides and build them
for item in slide_data:
    s = make_base_slide(item["num"], item["title"], item["presenter"], item["notes"])
    
    t = item["content_type"]
    if t == "title_slide":
        add_text_box(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.8),
                     item["header"], font_size=40, font_name='Cormorant Garamond',
                     color=IVORY, alignment=PP_ALIGN.CENTER, italic=True)
        add_gold_line(s, Inches(4.5), Inches(4.3), Inches(4.3))
        add_text_box(s, Inches(1), Inches(4.7), Inches(11.3), Inches(0.8),
                     item["sub"], font_size=16, font_name='Inter',
                     color=BODY_COLOR, alignment=PP_ALIGN.CENTER)
                     
    elif t == "metrics":
        add_text_box(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                     item["header"], font_size=36, font_name='Cormorant Garamond',
                     color=IVORY, italic=True)
        add_gold_line(s, Inches(0.8), Inches(2.2), Inches(3))
        
        # Grid arrangement
        for i, (name, role, hrs) in enumerate(item["metrics"]):
            col = i % 2
            row = i // 2
            x = Inches(0.8 + col * 5.8)
            y = Inches(2.6 + row * 2.0)
            
            # Card background
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.3), Inches(1.6))
            card.fill.solid()
            card.fill.fore_color.rgb = SURFACE
            card.line.color.rgb = MUTED
            card.line.width = Pt(0.5)
            
            # Text inside card
            add_text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(4.7), Inches(0.3),
                         role.upper(), font_size=10, font_name='JetBrains Mono', color=GOLD)
            add_text_box(s, x + Inches(0.3), y + Inches(0.55), Inches(4.7), Inches(0.4),
                         name, font_size=18, font_name='Cormorant Garamond', color=IVORY)
            add_text_box(s, x + Inches(0.3), y + Inches(1.05), Inches(4.7), Inches(0.35),
                         f"Zeitaufwand: {hrs}", font_size=12, font_name='Inter', color=BODY_COLOR)

    elif t == "two_col":
        add_text_box(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                     item["header"], font_size=36, font_name='Cormorant Garamond',
                     color=IVORY, italic=True)
        add_gold_line(s, Inches(0.8), Inches(2.2), Inches(3))
        
        # Left column
        add_text_box(s, Inches(0.8), Inches(2.6), Inches(5.3), Inches(0.4),
                     item["col1_title"], font_size=12, font_name='JetBrains Mono', color=GOLD, bold=True)
        add_text_box(s, Inches(0.8), Inches(3.2), Inches(5.3), Inches(3.2),
                     item["col1_text"], font_size=14, font_name='Inter', color=BODY_COLOR)
                     
        # Right column
        add_text_box(s, Inches(6.8), Inches(2.6), Inches(5.5), Inches(0.4),
                     item["col2_title"], font_size=12, font_name='JetBrains Mono', color=GOLD, bold=True)
        add_text_box(s, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.2),
                     item["col2_text"], font_size=14, font_name='Inter', color=BODY_COLOR)

    elif t == "three_col":
        add_text_box(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                     item["header"], font_size=36, font_name='Cormorant Garamond',
                     color=IVORY, italic=True)
        add_gold_line(s, Inches(0.8), Inches(2.2), Inches(3))
        
        cols = [item["col1"], item["col2"], item["col3"]]
        for i, (kicker, title, desc) in enumerate(cols):
            x = Inches(0.8 + i * 4.0)
            y = Inches(2.6)
            
            # Column Card
            card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(3.8))
            card.fill.solid()
            card.fill.fore_color.rgb = SURFACE
            card.line.color.rgb = MUTED
            card.line.width = Pt(0.5)
            
            add_text_box(s, x + Inches(0.25), y + Inches(0.3), Inches(3.2), Inches(0.35),
                         kicker, font_size=10, font_name='JetBrains Mono', color=GOLD)
            add_text_box(s, x + Inches(0.25), y + Inches(0.7), Inches(3.2), Inches(0.8),
                         title, font_size=18, font_name='Cormorant Garamond', color=IVORY, bold=True)
            add_text_box(s, x + Inches(0.25), y + Inches(1.6), Inches(3.2), Inches(2.0),
                         desc, font_size=12, font_name='Inter', color=BODY_COLOR)

# Save presentation
output_path = r"c:\Users\HP\Desktop\abschlussprojekt\04_praesentationen\Mirrou_Abschlusspraesentation_DCI.pptx"
prs.save(output_path)
print(f"PowerPoint Presentation successfully saved: {output_path}")
print(f"Total Slides: {len(prs.slides)}")
